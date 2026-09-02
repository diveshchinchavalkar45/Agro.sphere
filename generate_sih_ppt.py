from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation("SIH2026-IDEA-Presentation-Format.pptx")

# Delete slide 7 (instructions slide)
rId = prs.slides._sldIdLst[6].rId
prs.part.drop_rel(rId)
del prs.slides._sldIdLst[6]

# Helper to update team name oval on slides
team_name = "Agro-Sphere Innovators"
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame and "Your Team Name" in shape.text_frame.text:
            shape.text_frame.text = team_name
            for p in shape.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.bold = True

# ==========================================
# SLIDE 1: TITLE PAGE
# ==========================================
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame and "Problem Statement ID" in shape.text_frame.text:
        tf = shape.text_frame
        tf.clear()
        
        items = [
            ("Problem Statement ID: ", "SIH1645 (Smart Agriculture & Farm Connect)"),
            ("Problem Statement Title: ", "AI Mandi Price Intelligence, Quality Grading & Farm Logistics"),
            ("Theme: ", "Agriculture, FoodTech & Rural Development"),
            ("PS Category: ", "Software"),
            ("Team ID: ", "SIH2026-TEAM-AGRO"),
            ("Team Name: ", "Agro-Sphere Innovators"),
            ("Project / Idea: ", "Agro-Sphere (Multilingual Farmer Intelligence Platform)"),
            ("Live Prototype: ", "https://diveshchinchavalkar45.github.io/Agro.sphere/")
        ]
        
        for label, val in items:
            p = tf.add_paragraph()
            r1 = p.add_run()
            r1.text = label
            r1.font.bold = True
            r1.font.size = Pt(12)
            r1.font.color.rgb = RGBColor(23, 58, 37)
            
            r2 = p.add_run()
            r2.text = val
            r2.font.bold = False
            r2.font.size = Pt(12)
            r2.font.color.rgb = RGBColor(45, 45, 45)
            p.space_after = Pt(4)

# Add Logo to Slide 1
slide1.shapes.add_picture("assets/logo.png", Inches(8.3), Inches(4.3), width=Inches(1.6))

# ==========================================
# SLIDE 2: PROPOSED SOLUTION
# ==========================================
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame and "Proposed Solution" in shape.text_frame.text:
        tf = shape.text_frame
        tf.clear()
        
        lines = [
            ("Agro-Sphere: Unified Farm Intelligence & Direct Market Ecosystem", True, 13, RGBColor(23, 58, 37)),
            ("• Transparent Mandi Discovery: Live modal rates across APMC mandis (Lasalgaon, Nashik, Pune) with predictive 7-day selling recommendations.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Verified Quality Grading & Inspection: Standardized lot specs (moisture %, size 45-55mm, zero residue test, defect <2%) for Grade A, B & Premium produce.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Zero-Waste Fast Clearance: Emergency fast-sale channel for perishable surplus fruits/vegetables to eliminate post-harvest crop rotting.", False, 10.5, RGBColor(40, 40, 40)),
            ("• FPO Collective Pooling: Aggregates small farmer harvests into 100q+ bulk lots to secure institutional buyer pricing & reduce freight.", False, 10.5, RGBColor(40, 40, 40)),
            ("• End-to-End Live Logistics: Real-time truck tracking (MH-15), driver contact dispatch, weighbridge slot booking, and digital calendar milestones.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Multilingual Voice AI Sahayak: Voice-operated in Hindi, Marathi & English for non-literate farmers with 24x7 Toll-Free Helpline (1800-889-2476).", False, 10.5, RGBColor(40, 40, 40))
        ]
        
        for text, bold, size, color in lines:
            p = tf.add_paragraph()
            p.text = text
            p.font.bold = bold
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_after = Pt(3)

slide2.shapes.add_picture("products_modal_shot.png", Inches(6.8), Inches(1.8), width=Inches(3.0))

# ==========================================
# SLIDE 3: TECHNICAL APPROACH
# ==========================================
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.has_text_frame and "Technologies to be used" in shape.text_frame.text:
        tf = shape.text_frame
        tf.clear()
        
        lines = [
            ("Technical Architecture & Execution Methodology", True, 13, RGBColor(23, 58, 37)),
            ("• Modern Web & Mobile Stack: Responsive Glassmorphism architecture, CSS Grid/Flexbox, dynamic modular section routing.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Voice AI Engine: Web Speech API (STT & TTS) + Vapi.ai Telephony Integration supporting Hindi, Marathi, and Indian English.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Logistics & Scheduling Engine: Real-time route progression, live GPS coordinate simulator, and dynamic monthly milestone calendar.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Multilingual I18n System: Real-time DOM dictionary translation across 12 Indian regional languages.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Security & Settlement: PIN-verified weighbridge authorization and instant ledger balance tracking.", False, 10.5, RGBColor(40, 40, 40))
        ]
        
        for text, bold, size, color in lines:
            p = tf.add_paragraph()
            p.text = text
            p.font.bold = bold
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_after = Pt(2.5)

slide3.shapes.add_picture("architecture_flowchart.png", Inches(0.8), Inches(4.2), width=Inches(4.8))
slide3.shapes.add_picture("schedule_shot.png", Inches(5.8), Inches(4.2), width=Inches(3.9))

# ==========================================
# SLIDE 4: FEASIBILITY AND VIABILITY
# ==========================================
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if shape.has_text_frame and "Analysis of the feasibility" in shape.text_frame.text:
        tf = shape.text_frame
        tf.clear()
        
        lines = [
            ("Feasibility, Risk Analysis & Mitigation Strategies", True, 13, RGBColor(23, 58, 37)),
            ("• Technical Feasibility: Lightweight client architecture works on low-bandwidth 2G/3G rural networks with offline resilience.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Operational Feasibility: Multi-channel support (24x7 Toll-Free 1800-889-2476, WhatsApp Photo Grading +91 98200 45678, on-ground APMC officers).", False, 10.5, RGBColor(40, 40, 40)),
            ("• Financial Viability: Disintermediates exploitative middlemen (saving 15-20% margin for farmers); self-sustainable via minimal 0.5% buyer transaction fee.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Challenge: Low Digital Literacy among elderly farmers.", False, 10.5, RGBColor(180, 40, 40)),
            ("  ↳ Mitigation: Multilingual AI Voice Sahayak with one-touch voice chips and natural speech queries.", False, 10.5, RGBColor(23, 58, 37)),
            ("• Challenge: Weight & Quality Disputes at Mandi Yard.", False, 10.5, RGBColor(180, 40, 40)),
            ("  ↳ Mitigation: Digital weighbridge PIN verification receipt and local field officer dispute desk.", False, 10.5, RGBColor(23, 58, 37)),
            ("• Challenge: Spoilage of Perishable Produce.", False, 10.5, RGBColor(180, 40, 40)),
            ("  ↳ Mitigation: Priority Zero-Waste Fast Clearance tag and instant buyer flash-sale notifications.", False, 10.5, RGBColor(23, 58, 37))
        ]
        
        for text, bold, size, color in lines:
            p = tf.add_paragraph()
            p.text = text
            p.font.bold = bold
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_after = Pt(2)

# ==========================================
# SLIDE 5: IMPACT AND BENEFITS
# ==========================================
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.has_text_frame and "Potential impact" in shape.text_frame.text:
        tf = shape.text_frame
        tf.clear()
        
        lines = [
            ("Measurable Impact & Multi-Dimensional Benefits", True, 13, RGBColor(23, 58, 37)),
            ("• Economic Empowerment (+18% to +25% Income): Direct market access eliminates local trader cartels, delivering higher net realized price.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Zero Post-Harvest Food Wastage: Urgent clearance sell feature cuts agricultural perishable dumping by over 30%.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Farmer Collectives & FPO Growth: Collective pooling aggregates smallholder output, granting small farmers enterprise bargaining power.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Supply Chain Efficiency: Guaranteed transport dispatch timelines eliminate dead freight and multi-day mandi yard waiting times.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Financial Security & Escrow Trust: Instant digital bank settlement upon weighbridge slip verification prevents payment defaults.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Inclusive Accessibility: Voice-first multi-lingual interface bridges the digital divide for every Indian kisan.", False, 10.5, RGBColor(40, 40, 40))
        ]
        
        for text, bold, size, color in lines:
            p = tf.add_paragraph()
            p.text = text
            p.font.bold = bold
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_after = Pt(2.5)

slide5.shapes.add_picture("dashboard_bottom_shot.png", Inches(6.2), Inches(2.2), width=Inches(3.5))

# ==========================================
# SLIDE 6: RESEARCH AND REFERENCES
# ==========================================
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if shape.has_text_frame and "Details / Links" in shape.text_frame.text:
        tf = shape.text_frame
        tf.clear()
        
        lines = [
            ("Research Grounding, Live Deployment & References", True, 13, RGBColor(23, 58, 37)),
            ("• Live Deployed Web Prototype: https://diveshchinchavalkar45.github.io/Agro.sphere/", True, 10.5, RGBColor(33, 115, 70)),
            ("• Open Source Codebase: https://github.com/diveshchinchavalkar45/Agro.sphere", False, 10.5, RGBColor(40, 40, 40)),
            ("• Government Data Reference: Agmarknet (Directorate of Marketing & Inspection, Ministry of Agriculture & Farmers Welfare).", False, 10.5, RGBColor(40, 40, 40)),
            ("• Post-Harvest Loss Research: NITI Aayog Strategy for Doubling Farmers Income & Central Institute of Post-Harvest Engineering (CIPHET).", False, 10.5, RGBColor(40, 40, 40)),
            ("• Quality & Grading Standards: Bureau of Indian Standards (BIS) & APEDA Export Produce Specifications for Fruits & Vegetables.", False, 10.5, RGBColor(40, 40, 40)),
            ("• Voice AI Technology: Deepgram Multilingual Speech Recognition & Vapi.ai Telephony Voice Pipelines.", False, 10.5, RGBColor(40, 40, 40))
        ]
        
        for text, bold, size, color in lines:
            p = tf.add_paragraph()
            p.text = text
            p.font.bold = bold
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_after = Pt(3)

slide6.shapes.add_picture("prices_shot.png", Inches(5.8), Inches(3.2), width=Inches(3.8))

prs.save("Agro_Sphere_SIH2026_Idea_Presentation.pptx")
print("Saved final polished Agro_Sphere_SIH2026_Idea_Presentation.pptx!")
